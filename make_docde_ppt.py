# -*- coding: utf-8 -*-
# make_docde_ppt.py
# Tạo PowerPoint 15 slide giới thiệu “Đọc Dễ – Hỗ trợ trẻ rối loạn đọc”
# Yêu cầu: Python 3.9+ ; pip install --upgrade python-pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Bảng màu
INK = RGBColor(0x37, 0x41, 0x51)      # #374151
ACC = RGBColor(0x0e, 0xa5, 0xa0)      # #0ea5a0
MUT = RGBColor(0x51, 0x62, 0x6f)      # #51626f
PILL = RGBColor(0xF4, 0xF6, 0xFB)     # #f4f6fb

prs = Presentation()
# 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_title(target, text, color=INK, size=42, bold=True, align=PP_ALIGN.LEFT):
    """Đặt tiêu đề; target là Shape (có .text_frame) hoặc TextFrame."""
    tf = target.text_frame if hasattr(target, "text_frame") else target
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.bold = bold
    r.font.size = Pt(size)
    r.font.color.rgb = color
    p.alignment = align

def add_footer(slide, text="Đọc Dễ – ddh14.github.io/App_for_dyslexia_children/"):
    left = Inches(0.3)
    top = prs.slide_height - Inches(0.5)
    width = prs.slide_width - Inches(0.6)
    height = Inches(0.3)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.color.rgb = MUT

def add_header_bar(slide, color=ACC):
    # Thanh màu mảnh phía trên (không viền)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.18))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.width = 0

def add_pill(slide, left, top, text, fg=INK, bg=PILL):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.2), Inches(0.6))
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.color.rgb = RGBColor(0xE6, 0xE8, 0xEE)
    tx = box.text_frame
    tx.clear()
    p = tx.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(14)
    r.font.color.rgb = fg
    tx.margin_left = Inches(0.2)
    tx.margin_right = Inches(0.2)
    tx.margin_top = Inches(0.05)

def add_bullets(slide, items, left=Inches(0.8), top=Inches(2.0),
                width=None, height=None, size=24):
    if width is None:
        width = prs.slide_width - Inches(1.6)
    if height is None:
        height = prs.slide_height - top - Inches(1.3)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.clear()
    for i, txt in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = INK

def add_illus_block(slide, emoji="📊", caption="Minh hoạ"):
    # Khối minh hoạ bên phải
    w = Inches(4.2); h = Inches(3.2)
    left = prs.slide_width - w - Inches(0.6); top = Inches(2.0)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    shape.line.color.rgb = RGBColor(0xE6, 0xE8, 0xEE)

    # Emoji lớn
    box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.4), w - Inches(0.4), h - Inches(1.0))
    tf = box.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = emoji; r.font.size = Pt(80)

    # Chú thích
    cap_box = slide.shapes.add_textbox(left, top + h - Inches(0.7), w, Inches(0.6))
    cap = cap_box.text_frame; cap.text = caption
    for rr in cap.paragraphs[0].runs:
        rr.font.size = Pt(14); rr.font.color.rgb = MUT
    cap.paragraphs[0].alignment = PP_ALIGN.CENTER

def ensure_title_shape(slide):
    """Trả về shape tiêu đề; nếu layout không có title -> tạo textbox mới."""
    t = slide.shapes.title
    if t is None:
        t = slide.shapes.add_textbox(Inches(0.8), Inches(0.8),
                                     prs.slide_width - Inches(1.6), Inches(1.0))
    return t

def add_title_content_slide(title, bullets, emoji="🧩", cap="Minh hoạ"):
    layout = prs.slide_layouts[5]  # TITLE_ONLY
    slide = prs.slides.add_slide(layout)
    add_header_bar(slide)
    set_title(ensure_title_shape(slide), title, size=34)
    add_bullets(slide, bullets)
    add_illus_block(slide, emoji=emoji, caption=cap)
    add_footer(slide)
    return slide

# Slide 1 – Bìa
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_header_bar(slide)

# Tiêu đề lớn
title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5),
                                     prs.slide_width - Inches(1.6), Inches(1.6))
tf = title_box.text_frame; tf.clear()
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Đọc Dễ – Ứng dụng hỗ trợ trẻ rối loạn đọc"
r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = INK

# Phụ đề
sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(8.8), Inches(1.0))
sub = sub_box.text_frame; sub.text = "Structured Literacy • PA • Cards (SRS) • Luyện đọc • Dashboard"
for run in sub.paragraphs[0].runs:
    run.font.size = Pt(22); run.font.color.rgb = MUT

# Pills
add_pill(slide, Inches(0.8), Inches(3.8), "Voice UI 🔊")
add_pill(slide, Inches(3.1), Inches(3.8), "TTS & Ghi âm")
add_pill(slide, Inches(5.4), Inches(3.8), "Offline‑first")

# Link
link_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(9.5), Inches(0.6))
link = link_box.text_frame; link.text = "Demo: ddh14.github.io/App_for_dyslexia_children/"
for run in link.paragraphs[0].runs:
    run.font.size = Pt(16); run.font.color.rgb = ACC

add_illus_block(slide, emoji="📖🔊🃏", caption="PA • Cards • Reading")
add_footer(slide)

# Slide 2–15 (nội dung rút gọn từ báo cáo)
add_title_content_slide(
    "Vấn đề & bối cảnh",
    [
        "Dyslexia ước 5–10% HS; khó nhận âm–chữ, thanh điệu, tốc độ đọc.",
        "Thiếu công cụ tiếng Việt; phụ huynh khó duy trì phiên đều.",
        "Cần giải pháp: ngắn – đều – đúng nhóm lỗi; có dữ liệu tiến bộ."
    ],
    emoji="🧒📚", cap="Bối cảnh lớp học"
)
add_title_content_slide(
    "Mục tiêu của Đọc Dễ",
    [
        "Tự động hoá nhận diện–phát âm tiếng Việt có dấu.",
        "Giảm lỗi chính tả theo ngữ cảnh (s/x; ch/tr; ng/ngh; g/gh; c/k/qu; n/l…).",
        "Tăng WCPM & hiểu nội dung; theo dõi tiến bộ & gợi ý ưu tiên."
    ],
    emoji="🎯📈", cap="Đích đến"
)
add_title_content_slide(
    "Cách tiếp cận sư phạm",
    [
        "Structured Literacy + đa giác quan (Nghe–Nhìn–Nhấn–Nói).",
        "Chuỗi: PA → Cards (SRS) → Reading; phiên 4–7’/ngày.",
        "Quy trình: Cue → Nghe → Thực hiện → Phản hồi → Củng cố."
    ],
    emoji="➡️", cap="Quy trình 5 bước"
)
add_title_content_slide(
    "Kiến trúc & công nghệ",
    [
        "PWA offline‑first; Web/Android/iOS; Cordova/Capacitor.",
        "TTS vi‑VN; Voice UI; Recorder; Spotlight/Pacer/Echo.",
        "Dashboard, CSV/PDF; đồng bộ Apps Script (opt‑in)."
    ],
    emoji="🧱⚙️", cap="Kiến trúc"
)
add_title_content_slide(
    "PA – Segment (ghép mảnh)",
    [
        "Ví dụ: “tr–anh → tranh”, “gh–é → ghé”, “qu–ả → quả”.",
        "Nhấn quy tắc c/k/qu; ng/ngh; g/gh theo nguyên âm sau.",
        "Phản hồi tức thì; gợi ý ngắn theo mẫu chữ."
    ],
    emoji="🧩", cap="Ghép onset–rime"
)
add_title_content_slide(
    "PA – Tone & Pair",
    [
        "Tone: 6 thanh cùng base; mã màu nhất quán; cặp khó sắc/ngã, huyền/hỏi.",
        "Pair: cặp tối thiểu theo tag (s/x, ch/tr, ng/ngh, g/gh, c/k/qu, n/l, d/gi/r).",
        "Gợi ý tầng bậc: 🐢 nghe chậm → manh mối → đối chiếu."
    ],
    emoji="🎵⚖️", cap="Thanh & cặp tối thiểu"
)
add_title_content_slide(
    "Cards (SRS)",
    [
        "Thẻ: 1 tiếng có dấu + emoji; tô màu theo 6 thanh.",
        "SM‑2 rút gọn: Q=5/3/1 → khoảng cách 1–60 ngày; “đã vững” I≥14d.",
        "Vòng 60s tăng nhịp; filter tag; gợi ý tag yếu 7 ngày."
    ],
    emoji="🃏", cap="Ôn đúng lúc"
)
add_title_content_slide(
    "Reading – 5 bước & chỉ số",
    [
        "B1 Chọn • B2 Đọc (ghi âm) • B3 Đánh dấu • B4 Hỏi • B5 Tổng kết.",
        "Chỉ số: WCPM, Accuracy, lỗi theo loại; tự sửa theo cấu hình.",
        "Công cụ: Spotlight/Pacer/Echo; Replay nghe lại."
    ],
    emoji="📖⏱️", cap="Flow 1→5"
)
add_title_content_slide(
    "Game & Adaptive/Dashboard",
    [
        "Game: 60s; Tone/Tag mode; bóng đúng có vòng sáng; combo & điểm.",
        "Adaptive: ưu tiên tag yếu; trộn 70–80% mục tiêu + 20–30% dễ.",
        "Dashboard: WCPM/Accuracy; lỗi theo tag; due/mastered; gợi ý tuần."
    ],
    emoji="🎈📊", cap="Động lực & theo dõi"
)
add_title_content_slide(
    "Quyền riêng tư & lưu trữ",
    [
        "Thu tối thiểu – ẩn danh – do người dùng kiểm soát.",
        "Lưu localStorage/IndexedDB; đồng bộ HTTPS+SECRET (opt‑in).",
        "Unicode NFC; Intl.Segmenter('vi'); SW offline‑first."
    ],
    emoji="🛡️", cap="An toàn dữ liệu"
)
add_title_content_slide(
    "Hạn chế",
    [
        "Khác biệt TTS (đặc biệt hỏi/ngã/nặng).",
        "SRS dựa tự chấm có thiên lệch; cần giám sát nhẹ.",
        "Phương ngữ & phạm vi học liệu; quy mô mẫu còn nhỏ."
    ],
    emoji="⚠️", cap="Nhận diện thách thức"
)
add_title_content_slide(
    "Khả năng mở rộng",
    [
        "Mở rộng tag/level/phương ngữ; Reading 3–6; chiến lược đọc.",
        "Gói học liệu ký số; công cụ giáo viên; guided/repeated reading.",
        "Cải thiện TTS/contour; heatmap lỗi; phân tích tăng trưởng."
    ],
    emoji="🧭", cap="Roadmap kỹ thuật"
)
add_title_content_slide(
    "Đánh giá – phương án & công cụ",
    [
        "Đối tượng: Chuyên gia • Phụ huynh • Học sinh; thiết bị Android/iOS.",
        "Phiên 10–20 phút; 4 phiếu: CG, PH, HS, Quan sát; quy trình đồng thuận.",
        "Ẩn danh; không thu PII/audio mặc định; opt‑in đồng bộ."
    ],
    emoji="📝", cap="Thiết kế đánh giá"
)
add_title_content_slide(
    "Kết quả chính",
    [
        "Hoàn thành đo cao; thời lượng bám khuyến nghị; bấm nhầm thấp.",
        "Hứng thú vừa–cao; nhu cầu nghe chậm chủ yếu ở hỏi/ngã/nặng.",
        "Đồng thuận cao: sư phạm, thanh, UI, Spotlight/Pacer/Echo, hiển thị tiếng Việt."
    ],
    emoji="✅", cap="Tín hiệu tích cực"
)
add_title_content_slide(
    "Kết luận",
    [
        "Giải pháp khoa học–nhân văn, chi phí thấp, dễ nhân rộng.",
        "Kết nối lí thuyết đọc với Structured Literacy Việt hoá.",
        "Hướng tiếp: mở rộng học liệu/giọng TTS, công cụ GV, theo dõi dài hạn."
    ],
    emoji="🏁", cap="Thông điệp chốt"
)

prs.save("DocDe_GioiThieu.pptx")
print("Đã xuất: DocDe_GioiThieu.pptx")